#!/usr/bin/env python3
"""
东方商务社交空间 — 美团推流落地行动方案 PPT
东方美学风格 · 暗绿+金色系
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Color Palette ===
DARK_GREEN = RGBColor(0x1A, 0x3C, 0x2E)       # 主色 - 深墨绿
GOLD = RGBColor(0xC8, 0xA8, 0x5A)              # 金色
LIGHT_GOLD = RGBColor(0xE8, 0xD5, 0x9E)        # 浅金
CREAM = RGBColor(0xF5, 0xF0, 0xE6)             # 米白
DARK_BG = RGBColor(0x12, 0x2C, 0x22)           # 暗底
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xE8, 0xE0, 0xD0)        # 柔白
ACCENT_RED = RGBColor(0x8B, 0x3A, 0x2A)        # 朱红点缀
TEXT_DARK = RGBColor(0x2C, 0x2C, 0x2C)
TEXT_LIGHT = RGBColor(0xAA, 0xA0, 0x88)
SUBTLE_GREEN = RGBColor(0x2D, 0x5A, 0x45)       # 次绿

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=DARK_BG):
    """Fill slide background"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_LIGHT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=14, default_color=TEXT_LIGHT, line_spacing=1.5):
    """Add multi-line text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', default_size))
            p.font.color.rgb = item.get('color', default_color)
            p.font.bold = item.get('bold', False)
            p.font.name = item.get('font', 'Microsoft YaHei')
            p.space_after = Pt(item.get('space_after', 6))
        else:
            p.text = str(item)
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
            p.font.name = 'Microsoft YaHei'
            p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT
    return txBox

def add_accent_line(slide, left, top, width, color=GOLD, height=Pt(2)):
    return add_shape(slide, left, top, width, height, color)

def add_title_slide(title, subtitle="", slide_num=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)
    
    # Left accent bar
    add_shape(slide, Inches(0.8), Inches(0), Pt(4), H, GOLD)
    
    # Top gold line
    add_shape(slide, Inches(0), Inches(1.2), Inches(5), Pt(1), LIGHT_GOLD)
    
    # Title
    add_text_box(slide, Inches(1.2), Inches(1.6), Inches(10), Inches(1.2),
                 title, font_size=36, color=GOLD, bold=True)
    
    # Subtitle
    if subtitle:
        ac = add_accent_line(slide, Inches(1.2), Inches(2.9), Inches(2))
        add_text_box(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=SOFT_WHITE)
    
    # Bottom info
    add_text_box(slide, Inches(1.2), Inches(6.5), Inches(5), Inches(0.5),
                 f"成都高新区 · 招商大魔方", font_size=12, color=TEXT_LIGHT)
    add_text_box(slide, Inches(10), Inches(6.5), Inches(3), Inches(0.5),
                 f"2026.05", font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.RIGHT)
    
    # Decorative corner element
    add_shape(slide, Inches(12.5), Inches(0.5), Inches(0.3), Inches(0.3), GOLD)
    add_shape(slide, Inches(12.5), Inches(0.5), Inches(0.3), Pt(1), GOLD)
    
    return slide

def add_section_slide(section_num, section_title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    
    # Left gold bar
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, GOLD)
    
    # Section number
    add_text_box(slide, Inches(1), Inches(1.5), Inches(3), Inches(1),
                 f"0{section_num}", font_size=60, color=GOLD, bold=True)
    
    # Section title
    add_text_box(slide, Inches(1), Inches(3.0), Inches(10), Inches(1),
                 section_title, font_size=32, color=WHITE, bold=True)
    
    if subtitle:
        add_accent_line(slide, Inches(1), Inches(4.2), Inches(3), LIGHT_GOLD)
        add_text_box(slide, Inches(1), Inches(4.5), Inches(10), Inches(0.8),
                     subtitle, font_size=16, color=SOFT_WHITE)
    
    return slide

def add_content_slide(title, items, slide_num=None):
    """Standard content slide with left accent"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    
    # Top bar
    add_shape(slide, Inches(0), Inches(0), W, Pt(3), GOLD)
    
    # Left accent
    add_shape(slide, Inches(0.5), Inches(0.5), Pt(3), H - Inches(1), GOLD)
    
    # Title
    add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.7),
                 title, font_size=26, color=GOLD, bold=True)
    
    # Underline
    add_shape(slide, Inches(1), Inches(1.2), Inches(4), Pt(1), LIGHT_GOLD)
    
    # Items
    y = Inches(1.6)
    for item in items:
        if isinstance(item, str):
            add_text_box(slide, Inches(1.2), y, Inches(11), Inches(0.35),
                         f"▸ {item}", font_size=15, color=SOFT_WHITE)
            y += Inches(0.35)
        elif isinstance(item, dict):
            if 'title' in item:
                add_text_box(slide, Inches(1.2), y, Inches(11), Inches(0.35),
                             item['title'], font_size=17, color=GOLD, bold=True)
                y += Inches(0.35)
                if 'desc' in item:
                    add_text_box(slide, Inches(1.6), y, Inches(10.5), Inches(0.3),
                                 item['desc'], font_size=14, color=TEXT_LIGHT)
                    y += Inches(0.3)
            elif 'line' in item:
                add_shape(slide, Inches(1.2), y, Inches(10), Pt(0.5), SUBTLE_GREEN)
                y += Inches(0.2)
    
    return slide

def add_table_slide(title, headers, rows, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    
    # Top bar
    add_shape(slide, Inches(0), Inches(0), W, Pt(3), GOLD)
    
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 title, font_size=24, color=GOLD, bold=True)
    add_shape(slide, Inches(0.8), Inches(1.0), Inches(4), Pt(1), LIGHT_GOLD)
    
    # Table
    n_rows = len(rows) + 1
    n_cols = len(headers)
    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(11.5)
    height = Inches(0.4 * n_rows)
    
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    
    # Set column widths
    col_widths = [Inches(1.5) if i == 0 else Inches(2.5) for i in range(n_cols)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GOLD
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = DARK_BG
            p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3E, 0x30)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x16, 0x32, 0x26)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = SOFT_WHITE
                p.alignment = PP_ALIGN.CENTER
    
    if note:
        add_text_box(slide, Inches(0.8), Inches(6.5), Inches(10), Inches(0.5),
                     note, font_size=11, color=TEXT_LIGHT)
    
    return slide

# =============================================
# SLIDE 1: Cover
# =============================================
s1 = add_title_slide("东方商务社交空间", "美团推流 · 落地行动方案")

# =============================================
# SLIDE 2: 核心定位
# =============================================
s2 = add_content_slide("核心定位", [
    {'title': '成都高新区 · 东方商务第三空间', 'desc': '不是茶馆，是"空间时间售卖"的商务社交场所'},
    {'line': True},
    {'title': '地址', 'desc': '招商大魔方 · 257㎡ · 月固定成本35,000元'},
    {'line': True},
    {'title': '目标客群（按优先级）', 'desc': '① 高新区商务人群 ② 创业者/小团队 ③ 新中产女性 ④ 小型商务局'},
    {'line': True},
    {'title': '一句话定位', 'desc': '"安静、体面、有文化感、适合谈事的东方空间"'},
])

# =============================================
# SLIDE 3: 基建总览
# =============================================
s3 = add_section_slide(1, "第一阶段：店铺基建", "第1-3天 · 美团店铺重构")

# =============================================
# SLIDE 4: 店铺定位
# =============================================
s4 = add_content_slide("Step 1 · 店铺定位与命名", [
    {'title': '✅ 正确命名', 'desc': '「新中式商务茶空间」或「东方美学商务茶馆」'},
    {'line': True},
    {'title': '❌ 禁止使用', 'desc': '「茶馆」「喝茶」「品茶」—— 太低频、太传统'},
    {'line': True},
    {'title': '核心关键词（写进店铺介绍）', 'desc': '商务接待、会议、独立包间、新中式、东方美学、高新区'},
    {'line': True},
    {'title': '美团本质', 'desc': '不是团购平台，是成都本地流量搜索引擎。用户搜"商务茶空间"时你要排第一'},
])

# =============================================
# SLIDE 5: 首页图片
# =============================================
s5 = add_content_slide("Step 2 · 首页图片（决定生死）", [
    {'title': '前3张图的要求（美团算法核心看点击率）'},
    {'title': '第1张 — 大包间全景', 'desc': '光影+木结构+东方感，参考安缦酒店风格'},
    {'title': '第2张 — 商务洽谈场景', 'desc': '人在谈事，茶具作点缀，安静体面'},
    {'title': '第3张 — 氛围细节', 'desc': '局部留白、茶器、木纹、光影层次'},
    {'line': True},
    {'title': '❌ 禁忌', 'desc': '茶叶特写 / 空桌子 / 菜单截图 / 茶艺师特写'},
    {'title': '风格参考', 'desc': '安缦暗调 · 无印良品极简 · 野兽派东方 · 新中式酒店'},
])

# =============================================
# SLIDE 6: 标题
# =============================================
s6 = add_content_slide("Step 3 · 美团标题模板", [
    {'title': '模板一（商务首选）', 'desc': '【商务接待｜会议洽谈｜新中式茶空间】'},
    {'line': True},
    {'title': '模板二（年轻流量）', 'desc': '【高新区新中式茶馆｜适合聊天谈事｜独立包间】'},
    {'line': True},
    {'title': '模板三（高转化）', 'desc': '【可会议·可商务接待·可包间聚会的新中式茶空间】'},
    {'line': True},
    {'title': '大包改名（重要）', 'desc': '不要叫"大包间"，改叫"22人商务沙龙空间"'},
])

# =============================================
# SLIDE 7: SKU套餐
# =============================================
s7 = add_table_slide("Step 4 · 场景化SKU套餐设计", 
    ["套餐", "价格", "目标客群", "核心卖点"],
    [
        ["商务洽谈双人", "198元", "商务双人", "茶饮+茶点+2h卡座"],
        ["小型会议4人", "398元", "小团队", "包间+投屏+茶饮"],
        ["商务局12人", "888元", "创业局/私董会", "中包+茶饮+果盘"],
        ["高端沙龙22人", "1688~2688元", "企业活动", "全套商务服务"],
    ],
    "❌ 禁止只卖茶！卖的是「场景+空间」")

# =============================================
# SLIDE 8: 推广阶段
# =============================================
s8 = add_section_slide(2, "第二阶段：点金推广", "第4-30天 · 2000元/月预算精打细算")

# =============================================
# SLIDE 9: 预算
# =============================================
s9 = add_table_slide("Step 5 · 月预算结构（2000元）",
    ["模块", "预算", "占比"],
    [
        ["搜索推广", "1200元", "60%"],
        ["包间专项推广", "600元", "30%"],
        ["霸王餐/KOC", "200元", "10%"],
    ],
    "每日预算：工作日40-60元 | 周末80元")

# =============================================
# SLIDE 10: 关键词
# =============================================
s10 = add_content_slide("Step 6 · 搜索推广关键词清单", [
    {'title': '核心关键词（必投）'},
    '商务茶馆 · 茶室 · 包间茶馆 · 会议室',
    '高新区茶馆 · 大魔方茶馆 · 金融城茶馆',
    '洽谈 · 商务接待 · 适合谈事',
    {'line': True},
    {'title': '扩展关键词'},
    '下午茶 · 安静聊天 · 团建 · 小型会议',
    '新中式茶馆 · 围炉煮茶 · 私密空间',
])

# =============================================
# SLIDE 11: 时段策略
# =============================================
s11 = add_content_slide("Step 7 · 投放时段策略", [
    {'title': '只投两个核心时段'},
    {'title': '13:00 - 18:00', 'desc': '商务下午场（核心利润时段）'},
    {'title': '19:00 - 22:00', 'desc': '社交/商务局（高客单时段）'},
    {'line': True},
    {'title': '❌ 严禁全天投放', 'desc': '凌晨烧钱 + 引来低质流量 + 浪费预算'},
    {'line': True},
    {'title': '客群画像', 'desc': '不是学生，是商务人士、创业者、白领、小团队'},
])

# =============================================
# SLIDE 12: 三阶段投流
# =============================================
s12 = add_table_slide("Step 8 · 三阶段投流策略",
    ["阶段", "时间", "日预算", "核心任务"],
    [
        ["建立权重", "第1-7天", "40元", "搜索推广覆盖核心关键词"],
        ["定向人群", "第8-20天", "60元", "高消费+商务兴趣定向"],
        ["包间专项", "第21-30天", "80元", "会议包间独立推广计划"],
    ],
    "核心目标：不是订单量，是「包间预订量」")

# =============================================
# SLIDE 13: 内容阶段
# =============================================
s13 = add_section_slide(3, "第三阶段：内容与口碑", "持续运营 · KOC + 数据监控 + 多平台联动")

# =============================================
# SLIDE 14: 霸王餐
# =============================================
s14 = add_content_slide("Step 9 · 霸王餐/KOC运营", [
    {'title': '每月2-4场霸王餐体验'},
    '定价88元：茶饮 + 茶点 + 90分钟体验',
    '要求：必须上传照片 + 写真实体验',
    '不是亏钱，是做「高质量图片评价」',
    {'line': True},
    {'title': '美团算法权重', 'desc': '图片评价权重非常高 → 直接提升门店综合质量分'},
    {'line': True},
    {'title': 'KOC合作建议', 'desc': '找本地小红书博主 → 体验+拍照+发笔记 → 双平台曝光'},
])

# =============================================
# SLIDE 15: 数据看板
# =============================================
s15 = add_table_slide("Step 10 · 每周核心数据监控",
    ["指标", "目标值", "说明"],
    [
        ["点击率", ">8%", "美团首页决定生死"],
        ["入店率", ">15%", "图片质量决定"],
        ["收藏率", ">5%", "关系后续推荐权重"],
        ["包间咨询量", "每周增长", "核心KPI"],
        ["点评图片数", "持续增长", "影响综合评分"],
    ],
    "美团排序 = 出价 × 门店综合质量（含点击率/转化率/好评率/图片质量）")

# =============================================
# SLIDE 16: 多平台联动
# =============================================
s16 = add_content_slide("Step 11 · 多平台联动矩阵", [
    {'title': '美团', 'desc': '搜索流量截流 → 场景化SKU + 精准投放'},
    {'title': '小红书', 'desc': '审美种草 → "成都高新区适合谈事的东方空间"'},
    {'title': '大众点评', 'desc': '信任转化 → 高质量图文评价积累'},
    {'title': '视频号', 'desc': '熟人传播 → 空间光影短视频'},
    {'title': '抖音本地', 'desc': '扩大曝光 → 定向成都本地用户'},
    {'line': True},
    {'title': '小红书爆款方向', 'desc': '《成都终于有适合成年人聊天的地方了》'},
    '《高新区最安静的新中式茶空间》',
    '《比咖啡馆更适合谈合作的地方》',
])

# =============================================
# SLIDE 17: 转化阶段
# =============================================
s17 = add_section_slide(4, "第四阶段：线下转化与复购", "第2个月起 · 空间利用率是核心KPI")

# =============================================
# SLIDE 18: 营收拆解
# =============================================
s18 = add_table_slide("Step 12 · 月营收目标拆解（6.8~7.5万）",
    ["收入来源", "数量", "客单价", "月营收"],
    [
        ["大包（22人）", "10场", "2000元", "20,000元"],
        ["中包（10-12人）", "20场", "800元", "16,000元"],
        ["小包+卡座", "240桌", "80元", "19,200元"],
        ["茶叶零售", "-", "-", "5,000元"],
        ["企业活动", "2场", "1500元", "3,000元"],
    ],
    "盈亏平衡点：月营业额5万 | 目标区间：8万~12万/月")

# =============================================
# SLIDE 19: 活动圈层
# =============================================
s19 = add_content_slide("Step 13 · 固定主题活动圈层", [
    {'title': '每周三晚 — 创业者茶会', 'desc': '创业者圈层 → 高频复购 + 口碑传播'},
    {'title': '每周五晚 — 东方疗愈夜谈', 'desc': '女性情绪消费 → 小红书自发传播'},
    {'title': '周末 — 小型分享沙龙', 'desc': '社交人群 → 企业合作契机'},
    {'title': '工作日下午 — 联合办公茶空间', 'desc': '白领/自由职业 → 工作日填充'},
    {'line': True},
    {'title': '核心逻辑', 'desc': '一场优质活动 = 20~50个精准用户。自然散客一天才几桌'},
    {'title': '企业合作方向', 'desc': '联系创业社群/女性社群/读书会/心理疗愈/AI分享会 → 免费空间置换流量'},
])

# =============================================
# SLIDE 20: 会员体系
# =============================================
s20 = add_content_slide("Step 14 · 东方会客厅会员", [
    {'title': '年费：1,999元/年'},
    {'line': True},
    {'title': '权益包'},
    '· 包间折扣（8折）',
    '· 专属茶品（每月一款限定）',
    '· 优先预定（无需排队）',
    '· 茶会活动（免费参与）',
    '· 企业会议优惠（免费试用投屏）',
    {'line': True},
    {'title': '❌ 不要做低价充值', 'desc': '"充值500送50"太low，配不上东方商务调性'},
])

# =============================================
# SLIDE 21: 避坑指南
# =============================================
s21 = add_section_slide(5, "第五阶段：避坑指南", "绝对不能犯的7个错误")

# =============================================
# SLIDE 22: 避坑内容
# =============================================
s22 = add_content_slide("❌ 绝对不能做的事", [
    {'title': '① 低价团购', 'desc': '吸引打牌+久坐+低消费人群，毁掉空间调性'},
    {'title': '② "9.9喝茶"', 'desc': '直接拉低品牌价值，永远别碰'},
    {'title': '③ 全天候投流', 'desc': '凌晨也在烧钱，引来低质流量'},
    {'title': '④ 变成棋牌室', 'desc': '高新区年轻人不会来，商业模式直接死亡'},
    {'title': '⑤ 只讲茶文化', 'desc': '用户不关心。买的是情绪+社交+空间价值'},
    {'title': '⑥ 只卖茶', 'desc': '坪效撑不起257㎡，空间时间才是核心商品'},
    {'title': '⑦ 疯狂买曝光', 'desc': '美团不是抖音！转化率比曝光重要一百倍'},
])

# =============================================
# SLIDE 23: 核心原则
# =============================================
s23 = add_content_slide("✅ 核心原则与最终定位", [
    {'title': '一句话最终定位', 'desc': '"成都高新区的东方商务第三空间"'},
    {'line': True},
    {'title': '核心商业模式'},
    '卖空间 → 不卖茶（包间费是第一利润）',
    '卖关系、卖氛围、卖时间、卖身份感',
    '70%精力做B端，30%精力做C端',
    {'line': True},
    {'title': '美团角色', 'desc': '不是主要盈利来源，是"搜索截流工具"。真正成交在线下'},
    {'line': True},
    {'title': '真正卖什么', 'desc': '"适合成年人谈事的地方"——这个需求在成都高新区非常大'},
])

# =============================================
# SLIDE 24: 执行时间线
# =============================================
s24 = add_table_slide("完整执行时间线",
    ["阶段", "时间", "核心任务", "预算"],
    [
        ["基建期", "第1-3天", "定位·图片·标题·SKU", "-"],
        ["推广期", "第4-7天", "搜索推广·关键词·时段", "40元/天"],
        ["优化期", "第8-20天", "定向人群·提升转化", "60元/天"],
        ["深耕期", "第21-30天", "包间专项·KOC·数据监控", "80元/天"],
        ["稳定期", "第2个月起", "活动圈层·会员·B端", "2000元/月"],
    ],
    "目标：月营收8~12万 | 包间利用率持续提升")

# =============================================
# SLIDE 25: 结束
# =============================================
s25 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s25, DARK_BG)

# Decorative elements
add_shape(s25, Inches(0), Inches(0), Inches(0.15), H, GOLD)
add_shape(s25, Inches(5), Inches(2.5), Inches(3), Pt(2), GOLD)

add_text_box(s25, Inches(2), Inches(2.8), Inches(9), Inches(1.2),
             "东方商务社交空间", font_size=40, color=GOLD, bold=True)
add_text_box(s25, Inches(2), Inches(4.0), Inches(9), Inches(0.8),
             "成都高新区 · 招商大魔方", font_size=20, color=SOFT_WHITE)
add_text_box(s25, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
             "从卖茶到卖空间 · 从流量到留量", font_size=14, color=TEXT_LIGHT)

# Decorative corner
add_shape(s25, Inches(12), Inches(6.5), Inches(0.8), Inches(0.5), GOLD)

# Save
output_path = "/home/sayelf/.openclaw/workspace/output/meituan-action-plan/美团推流行动方案.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"File size: {os.path.getsize(output_path) / 1024:.1f} KB")
print(f"Total slides: {len(prs.slides)}")
