#!/usr/bin/env python3
"""
东欧产品渠道落地完整方案 — PPT
工业B2B风格 · 深蓝+橙色系
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Color Palette ===
DARK_BLUE = RGBColor(0x0D, 0x1B, 0x2A)      # 主色
ACCENT_ORANGE = RGBColor(0xF0, 0x7A, 0x1F)    # 橙色高亮
LIGHT_BLUE = RGBColor(0x1B, 0x3A, 0x5C)       # 次蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xE0, 0xE8, 0xF0)
LIGHT_GRAY = RGBColor(0x88, 0x99, 0xAA)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)
DARK_CARD = RGBColor(0x15, 0x28, 0x3D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width; H = prs.slide_height

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def add_tb(slide, l, t, w, h, text, sz=18, color=SOFT_WHITE, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def add_multi(slide, l, t, w, h, items, default_sz=14, default_c=SOFT_WHITE):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, dict):
            p.text = item.get('text',''); p.font.size = Pt(item.get('size', default_sz))
            p.font.color.rgb = item.get('color', default_c); p.font.bold = item.get('bold', False)
            p.font.name = item.get('font','Microsoft YaHei'); p.space_after = Pt(item.get('space', 4))
        else:
            p.text = str(item); p.font.size = Pt(default_sz); p.font.color.rgb = default_c
            p.font.name = 'Microsoft YaHei'; p.space_after = Pt(4)
        p.alignment = PP_ALIGN.LEFT
    return tb

def title_slide(main_title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_shape(s, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)
    add_shape(s, Inches(0.8), Inches(1.5), Inches(5), Pt(3), ACCENT_ORANGE)
    add_tb(s, Inches(1), Inches(1.8), Inches(10), Inches(1), main_title, 36, WHITE, True)
    if sub: add_tb(s, Inches(1), Inches(3.2), Inches(10), Inches(0.6), sub, 16, LIGHT_GRAY)
    add_tb(s, Inches(1), Inches(6.5), Inches(6), Inches(0.4), "2026-05-19 · 数据源: Trade.gov / EU TED", 11, LIGHT_GRAY)
    return s

def section_slide(num, title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_shape(s, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)
    add_tb(s, Inches(1), Inches(1.5), Inches(2), Inches(0.8), f"0{num}", 48, ACCENT_ORANGE, True)
    add_tb(s, Inches(1), Inches(3.0), Inches(10), Inches(0.8), title, 30, WHITE, True)
    if sub: add_tb(s, Inches(1), Inches(4.0), Inches(10), Inches(0.5), sub, 14, LIGHT_GRAY)
    return s

def content_slide(title, items):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_shape(s, Inches(0), Inches(0), W, Pt(3), ACCENT_ORANGE)
    add_shape(s, Inches(0.5), Inches(0.5), Pt(3), H-Inches(1), ACCENT_ORANGE)
    add_tb(s, Inches(1), Inches(0.5), Inches(10), Inches(0.6), title, 24, WHITE, True)
    add_shape(s, Inches(1), Inches(1.1), Inches(4), Pt(1), ACCENT_ORANGE)
    y = Inches(1.5)
    for item in items:
        if isinstance(item, str):
            add_tb(s, Inches(1.2), y, Inches(11), Inches(0.35), f"▸ {item}", 14, SOFT_WHITE); y += Inches(0.35)
        elif isinstance(item, dict):
            if 'title' in item:
                add_tb(s, Inches(1.2), y, Inches(11), Inches(0.35), item['title'], 16, ACCENT_ORANGE, True)
                y += Inches(0.35)
                if 'desc' in item:
                    add_tb(s, Inches(1.6), y, Inches(10.5), Inches(0.3), item['desc'], 13, LIGHT_GRAY); y += Inches(0.3)
            elif 'line' in item:
                add_shape(s, Inches(1.2), y, Inches(10), Pt(0.5), LIGHT_BLUE); y += Inches(0.2)
            elif 'tag' in item:
                add_tb(s, Inches(1.2), y, Inches(1), Inches(0.3), item['tag'], 11, ACCENT_ORANGE, True)
                add_tb(s, Inches(2.4), y, Inches(10), Inches(0.3), item['text'], 13, SOFT_WHITE); y += Inches(0.3)
    return s

def table_slide(title, headers, rows, note=""):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_shape(s, Inches(0), Inches(0), W, Pt(3), ACCENT_ORANGE)
    add_tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), title, 22, WHITE, True)
    add_shape(s, Inches(0.8), Inches(1.0), Inches(4), Pt(1), ACCENT_ORANGE)
    n_r = len(rows)+1; n_c = len(headers)
    ts = s.shapes.add_table(n_r, n_c, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4*n_r))
    t = ts.table
    for i, h in enumerate(headers):
        c = t.cell(0,i); c.text = h; c.fill.solid(); c.fill.fore_color.rgb = ACCENT_ORANGE
        for p in c.text_frame.paragraphs: p.font.size=Pt(12); p.font.bold=True; p.font.color.rgb=DARK_BLUE; p.alignment=PP_ALIGN.CENTER
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = t.cell(ri+1, ci); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = DARK_CARD if ri%2==0 else RGBColor(0x12,0x25,0x38)
            for p in c.text_frame.paragraphs: p.font.size=Pt(11); p.font.color.rgb=SOFT_WHITE; p.alignment=PP_ALIGN.CENTER
    if note: add_tb(s, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4), note, 10, LIGHT_GRAY)
    return s

# ===== SLIDES =====

# S1: Cover
title_slide("东欧产品渠道落地完整方案", "钢结构房屋 · 变压器 · 移动电源 · 通用发动机")

# S2: Channel Matrix
table_slide("产品-渠道矩阵",
    ["产品", "HS编码", "分类", "渠道模式", "优先级"],
    [
        ["钢结构折叠集成房屋", "7308.90", "B2B工业", "招标+展会+分销商", "P0"],
        ["变压器", "8504.21", "B2B工业", "招标+能源EPC+展会", "P0"],
        ["移动电源", "8507.60", "B2B/B2C", "电商平台+分销商", "P1"],
        ["通用发动机", "8408.20", "B2B工业", "招标+工业分销+展会", "P1"],
    ],
    "数据源：Trade.gov Country Commercial Guides")

# S3: Section - EU TED
section_slide(1, "EU TED — 欧盟官方招标数据库", "ted.europa.eu · 完全免费 · 覆盖27国")

# S4: TED Details
content_slide("EU TED 招标系统详解", [
    {'tag': '费用', 'text': '完全免费（浏览+搜索+RSS订阅）'},
    {'tag': '规模', 'text': '年发布20万+招标公告，覆盖欧盟27国+EEA'},
    {'tag': '适用', 'text': '钢结构房屋（基建营房）、变压器（电网升级）、发动机（应急电源）'},
    {'line': True},
    {'title': '搜索关键词配置'},
    '钢结构: "steel structure" "modular building" "foldable house"',
    '变压器: "transformer" "power transformer" "electrical substation"',
    '发动机: "engine" "diesel generator" "power generator"',
    {'line': True},
    {'tag': 'RSS', 'text': '支持XML/RSS订阅，可自动化监控新招标'},
])

# S5: Country Procurement Platforms
table_slide("各国政府采购平台（全部免费）",
    ["国家", "平台", "注册", "适合产品"],
    [
        ["🇵🇱 波兰", "Platforma Zakupowa", "免费", "全部工业品"],
        ["🇷🇴 罗马尼亚", "SEAP (e-licitatie.ro)", "免费", "全部工业品"],
        ["🇨🇿 捷克", "NEN (nen.nipez.cz)", "免费", "全部工业品"],
        ["🇭🇺 匈牙利", "EKR (ekr.gov.hu)", "免费", "全部工业品"],
        ["🇺🇦 乌克兰", "ProZorro (prozorro.gov.ua)", "免费", "全部工业品·重建核心"],
    ],
    "ProZorro: 乌克兰战后重建采购全面启动，日发布500+招标")

# S6: Section - B2C Platforms
section_slide(2, "电商平台渠道", "Allegro · eMAG · Prom.ua")

# S7: Allegro
content_slide("Allegro.pl — 波兰#1 电商平台（移动电源最佳渠道）", [
    {'tag': '流量', 'text': '月访问2亿+ | 年GMV ~$120亿 | 1400万+活跃买家'},
    {'tag': '费用', 'text': '佣金8-15%（按品类）| 国际卖家2023年开放'},
    {'line': True},
    {'title': '入驻步骤：'},
    '① 注册商业账户 → allegro.pl/business',
    '② 提供公司注册证明',
    '③ 申请波兰VAT税号',
    '④ 波兰仓储或Allegro Smart! FBA',
    '⑤ 波兰语产品上架',
    {'line': True},
    {'tag': '适合', 'text': '移动电源 ✅ | 钢结构/变压器/发动机 ❌'},
])

# S8: eMAG + Prom.ua
content_slide("eMAG.ro & Prom.ua — 罗马尼亚 & 乌克兰", [
    {'title': 'eMAG.ro（罗马尼亚#1）'},
    '覆盖罗马尼亚、保加利亚、匈牙利 | 年GMV ~$30亿',
    '品类优势：电子、家电、DIY',
    '✅ 移动电源适合 | ❌ 工业品不适合',
    {'line': True},
    {'title': 'Prom.ua（乌克兰#1）'},
    'B2B/B2C混合 | 月访问1亿+ | 支持跨境卖家',
    '✅ 移动电源、小型发动机 | ⚠️ 钢结构房屋B2B板块可能',
    {'line': True},
    {'tag': '特别', 'text': '乌克兰战后重建 → 钢结构房屋+发电设备是刚需'},
])

# S9: Section - B2B
section_slide(3, "B2B平台 + 行业展会", "Alibaba · Europages · Bauma · 各国重建展")

# S10: B2B Platforms
table_slide("B2B跨境电商平台对比",
    ["平台", "优势", "劣势", "适合产品"],
    [
        ["Alibaba.com", "全球买家，工业品成熟", "竞争激烈", "全部4类"],
        ["Made-in-China.com", "中国供应商主场", "流量低于阿里", "全部4类"],
        ["Europages.com", "欧洲B2B目录，免费", "被动等待询盘", "全部4类"],
        ["GlobalSources.com", "中高端买家", "需要付费", "高价值产品"],
    ],
    "建议：优先Alibaba.com + Europages免费版")

# S11: Trade Fairs
table_slide("2026-2027 行业展会时间表",
    ["展会", "地点", "时间", "适合产品"],
    [
        ["Green Energy Poland", "🇵🇱 波兹南", "2026年4月", "变压器/移动电源"],
        ["ENERGIA Energy Fair", "🇷🇴 布加勒斯特", "2026年5月", "变压器/储能"],
        ["ReBuild Ukraine", "🇵🇱 华沙", "2026年下半年", "钢结构房屋/发电设备"],
        ["MSZ Energy", "🇵🇱 凯尔采", "2026年9月", "变压器/储能"],
        ["FOR ENERGY", "🇨🇿 俄斯特拉发", "2026年10月", "变压器/发动机"],
        ["Bauma", "🇩🇪 慕尼黑", "2027年4月", "钢结构/发动机"],
    ],
    "数据源：Trade.gov 展会数据库 + 各展会官网")

# S12: Section - Execution
section_slide(4, "按产品执行路线图", "分产品拆解落地步骤")

# S13: Steel Structure
content_slide("产品1：钢结构折叠集成房屋（HS 7308.90）", [
    {'tag': 'P0', 'text': 'EU TED 搜索 "modular building" 招标 → 本周'},
    {'tag': 'P0', 'text': 'ProZorro 注册+监控乌克兰重建招标 → 本周'},
    {'tag': 'P0', 'text': '波兰 Platforma Zakupowa 注册+搜索 → 本周'},
    {'tag': 'P1', 'text': 'ReBuild Ukraine 展会预定 → 2026'},
    {'tag': 'P1', 'text': 'Europages 免费注册企业主页 → 本周'},
    {'tag': 'P2', 'text': '罗马尼亚 SEAP 注册+搜索'},
    {'line': True},
    {'title': '关键需求'},
    'CE标志（EN 1090钢结构认证）| 波兰语技术文档 | 本地合作伙伴',
    '波兰$250B能源转型 → 基建配套需求巨大',
])

# S14: Transformers
content_slide("产品2：变压器（HS 8504.21）", [
    {'tag': 'P0', 'text': 'EU TED 搜索 "transformer" "substation" → 本周'},
    {'tag': 'P0', 'text': 'Trade.gov Poland Energy 数据已拉取'},
    {'tag': 'P0', 'text': 'Trade.gov Romania Energy 数据已拉取'},
    {'tag': 'P1', 'text': '各国政府采购平台注册 → 本周'},
    {'tag': 'P1', 'text': 'MSZ Energy展会定展位 → 9月'},
    {'tag': 'P2', 'text': 'Alibaba.com 页面优化'},
    {'line': True},
    {'title': '市场机会（数据源：Trade.gov）'},
    '🇵🇱 波兰：煤电70%→56% | 电网改造大量变压器需求',
    '🇷🇴 罗马尼亚：储能仅16MW | 2032年退煤',
  ])

# S15: Power Banks
content_slide("产品3：移动电源（HS 8507.60）", [
    {'tag': 'P0', 'text': 'Allegro.pl 入驻+产品上架 → 本月'},
    {'tag': 'P0', 'text': 'eMAG.ro 评估入驻条件 → 本月'},
    {'tag': 'P1', 'text': 'Prom.ua 入驻乌克兰市场 → 本月'},
    {'tag': 'P1', 'text': 'EU TED 搜索应急采购（少量）'},
    {'line': True},
    {'title': '认证要求'},
    '波兰/欧盟: CE标志 + EN 62619 + WEEE电池回收注册',
    '乌克兰: UkrSEPRO 本地化认证',
    {'line': True},
    {'title': '市场机会（数据源：Trade.gov）'},
    '🇷🇴 罗马尼亚储能仅16MW → 户用储能从0爆发',
    '🇺🇦 乌克兰战后电力不稳定 → 移动电源刚需',
])

# S16: Engines
content_slide("产品4：通用发动机（HS 8408.20）", [
    {'tag': 'P0', 'text': 'EU TED 搜索 "engine" "generator" → 本周'},
    {'tag': 'P0', 'text': '各国政府采购 搜索应急/备用电源招标'},
    {'tag': 'P1', 'text': 'Alibaba.com B2B页面优化'},
    {'tag': 'P1', 'text': 'Bauma/行业展会参展或观众'},
    {'tag': 'P2', 'text': 'Europages 免费注册'},
])

# S17: Timeline
table_slide("执行时间线（3个月）",
    ["阶段", "时间", "核心动作"],
    [
        ["启动期", "Month1 W1", "EU TED注册+关键词配置"],
        ["启动期", "Month1 W2", "各国采购平台注册（波/罗/捷/匈/乌）"],
        ["启动期", "Month1 W3", "Allegro入驻准备（移动电源）"],
        ["启动期", "Month1 W4", "Europages + Alibaba页面优化"],
        ["增长期", "Month2 W1", "Allegro产品上架"],
        ["增长期", "Month2 W2", "Trade.gov定时抓取（cron）"],
        ["增长期", "Month2 W3", "展会注册（Green Energy/ReBuild）"],
        ["增长期", "Month2 W4", "EU TED监控脚本部署"],
        ["深耕期", "Month3 W1", "Allegro首月数据复盘"],
        ["深耕期", "Month3 W2", "联系波兰/罗马尼亚本地代理"],
        ["深耕期", "Month3 W3", "ProZorro乌克兰重建招标入场"],
        ["深耕期", "Month3 W4", "季度复盘+策略调整"],
    ],
    "所有渠道工具均为免费开源")

# S18: Free Tools
table_slide("免费工具清单",
    ["工具", "用途", "费用"],
    [
        ["ted.europa.eu", "欧盟招标搜索+RSS订阅", "免费"],
        ["prozorro.gov.ua", "乌克兰重建招标", "免费"],
        ["Platforma Zakupowa", "波兰政府采购", "免费"],
        ["SEAP (e-licitatie.ro)", "罗马尼亚政府采购", "免费"],
        ["Trade.gov", "市场情报+国别指南", "免费"],
        ["Europages", "欧洲B2B目录", "基础免费"],
        ["EU Access2Markets", "关税/规则查询", "免费"],
    ],
    "全部免费开源，无需任何付费API")

# S19: Data Sources
table_slide("数据来源",
    ["来源", "类型", "链接"],
    [
        ["Trade.gov — Poland Energy", "政府国别指南", "trade.gov/country-commercial-guides/poland-energy"],
        ["Trade.gov — Romania Energy", "政府国别指南", "trade.gov/country-commercial-guides/romania-energy"],
        ["EU TED", "欧盟招标数据库", "ted.europa.eu"],
        ["ProZorro", "乌克兰招标", "prozorro.gov.ua"],
        ["Allegro Business", "波兰电商", "allegro.pl/business"],
        ["eMAG", "罗马尼亚电商", "emag.ro"],
        ["Prom.ua", "乌克兰电商", "prom.ua"],
        ["Bauma", "工程机械展", "bauma.de"],
    ],
    "All sources are free and publicly accessible")

# S20: End
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
add_shape(s, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)
add_shape(s, Inches(5), Inches(2.8), Inches(3), Pt(2), ACCENT_ORANGE)
add_tb(s, Inches(2), Inches(3.2), Inches(9), Inches(1), "东欧产品渠道落地", 36, WHITE, True)
add_tb(s, Inches(2), Inches(4.3), Inches(9), Inches(0.6), "钢结构房屋 · 变压器 · 移动电源 · 通用发动机", 16, LIGHT_GRAY)
add_tb(s, Inches(2), Inches(5.5), Inches(9), Inches(0.5), "数据源: Trade.gov / EU TED / 各国政府采购平台", 12, LIGHT_GRAY)

# Save
out_dir = "/home/sayelf/.openclaw/workspace/output/geo-optimization"
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "东欧渠道落地方案.pptx")
prs.save(out_path)
print(f"✅ PPT: {out_path} ({os.path.getsize(out_path)/1024:.0f}KB, {len(prs.slides)} slides)")
