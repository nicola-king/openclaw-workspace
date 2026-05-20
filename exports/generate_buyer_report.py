#!/usr/bin/env python3
"""
买家情报 PDF/PPT 生成器
从 markdown 报告生成 PDF (WeasyPrint) + PPTX (python-pptx)
"""
import os
import sys
import json
import subprocess
from datetime import datetime
from pathlib import Path

sys.path.insert(0, os.path.expanduser("~/.openclaw/workspace/skills/cross-border-trade-agent/modules"))

MARKDOWN_FILE = os.path.expanduser("~/.openclaw/workspace/exports/buyer_intel_au_nz.md")
OUTPUT_DIR = os.path.expanduser("~/.openclaw/workspace/exports")
PDF_OUTPUT = os.path.join(OUTPUT_DIR, "买家情报报告_AU_NZ_20260520.pdf")
PPT_OUTPUT = os.path.join(OUTPUT_DIR, "买家情报报告_AU_NZ_20260520.pptx")
HTML_OUTPUT = os.path.join(OUTPUT_DIR, "buyer_intel_report_temp.html")


def md_to_html():
    """将 Markdown 转为美化 HTML"""
    with open(MARKDOWN_FILE) as f:
        md = f.read()

    html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<style>
@page {{
    size: A4;
    margin: 20mm 25mm;
}}
@font-face {{
    font-family: 'Noto Sans CJK';
    src: url('file:///usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc') format('truetype');
}}
body {{
    font-family: 'Noto Sans CJK', 'Helvetica Neue', Arial, sans-serif;
    color: #1a1a2e;
    line-height: 1.8;
    font-size: 10.5pt;
}}
h1 {{
    color: #16213e;
    border-bottom: 3px solid #0f3460;
    padding-bottom: 8px;
    font-size: 22pt;
    margin-top: 0;
}}
h2 {{
    color: #0f3460;
    border-bottom: 1px solid #e94560;
    padding-bottom: 4px;
    font-size: 16pt;
    margin-top: 28px;
}}
h3 {{
    color: #e94560;
    font-size: 13pt;
    margin-top: 18px;
}}
table {{
    width: 100%;
    border-collapse: collapse;
    margin: 12px 0;
    font-size: 9.5pt;
}}
th {{
    background: #16213e;
    color: white;
    padding: 6px 10px;
    text-align: left;
    font-weight: 600;
}}
td {{
    padding: 5px 10px;
    border-bottom: 1px solid #ddd;
}}
tr:nth-child(even) td {{
    background: #f8f9fa;
}}
hr {{
    border: none;
    border-top: 1px solid #e94560;
    margin: 20px 0;
}}
a {{
    color: #0f3460;
    text-decoration: none;
}}
.cover {{
    text-align: center;
    padding-top: 120px;
    page-break-after: always;
}}
.cover h1 {{
    font-size: 28pt;
    border: none;
    color: #16213e;
}}
.cover .subtitle {{
    font-size: 14pt;
    color: #e94560;
    margin-top: 16px;
}}
.cover .meta {{
    font-size: 10pt;
    color: #666;
    margin-top: 40px;
}}
.tag {{
    display: inline-block;
    background: #e94560;
    color: white;
    padding: 2px 8px;
    border-radius: 3px;
    font-size: 8pt;
    margin: 2px;
}}
.section-icon {{
    font-size: 18pt;
    margin-right: 8px;
}}
.priority-high {{
    background: #e94560; color: white; padding: 1px 6px; border-radius: 3px; font-size: 8pt;
}}
.priority-mid {{
    background: #f5a623; color: white; padding: 1px 6px; border-radius: 3px; font-size: 8pt;
}}
.priority-low {{
    background: #7ed321; color: white; padding: 1px 6px; border-radius: 3px; font-size: 8pt;
}}
blockquote {{
    border-left: 4px solid #0f3460;
    padding: 10px 16px;
    margin: 12px 0;
    background: #f0f2f5;
    font-size: 9pt;
}}
.footer {{
    text-align: center;
    font-size: 8pt;
    color: #999;
    margin-top: 40px;
    border-top: 1px solid #ddd;
    padding-top: 12px;
}}
</style>
</head>
<body>
"""
    # Convert markdown to simple HTML
    html += markdown_to_simple_html(md)
    html += '<div class="footer">'
    html += f'<p>太一跨境贸易 Agent · 买家情报引擎 · 生成日期：{datetime.now().strftime("%Y-%m-%d")}</p>'
    html += '</div></body></html>'
    return html


def markdown_to_simple_html(md):
    """Basic markdown to HTML conversion for WeasyPrint"""
    lines = md.split('\n')
    html_parts = []
    i = 0
    in_table = False

    while i < len(lines):
        line = lines[i]

        # Cover section
        if line.strip().startswith('# 终端买家情报报告'):
            html_parts.append('<div class="cover">')
            html_parts.append(f'<h1>终端买家情报报告</h1>')
            continue
        if line.strip().startswith('## 澳大利亚 & 新西兰'):
            html_parts.append(f'<p class="subtitle">澳大利亚 &amp; 新西兰 — 钢结构折叠集成房屋 / 模块化建筑</p>')
            html_parts.append(f'<p class="meta">生成日期：2026-05-20 | 数据来源：太一买家情报引擎<br>本报告为终端购买方情报，仅供参考</p>')
            html_parts.append('</div>')
            i += 1
            continue

        # Headings
        if line.startswith('### '):
            html_parts.append(f'<h3>{line[4:]}</h3>')
        elif line.startswith('## '):
            html_parts.append(f'<h2>{line[3:]}</h2>')
        elif line.startswith('# '):
            html_parts.append(f'<h1>{line[2:]}</h1>')

        # Tables
        elif '|' in line and line.strip().startswith('|'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            # Check if it's a separator row
            if all(c.replace('-','').replace(':','').strip() == '' for c in cells):
                in_table = True
            else:
                if not in_table or (i > 0 and lines[i-1].strip().startswith('|') and all(c.replace('-','').replace(':','').strip() == '' for c in [cc.strip() for cc in lines[i-1].split('|')[1:-1]])):
                    # Header row - start table
                    html_parts.append('<table><thead><tr>')
                    for c in cells:
                        html_parts.append(f'<th>{c}</th>')
                    html_parts.append('</tr></thead><tbody>')
                else:
                    html_parts.append('<tr>')
                    for c in cells:
                        # Handle links
                        c = c.replace('🔗 ', '<span style="font-size:10pt">🔗 </span>')
                        html_parts.append(f'<td>{c}</td>')
                    html_parts.append('</tr>')

        # End table on blank line
        elif in_table and not line.strip():
            html_parts.append('</tbody></table>')
            in_table = False

        # HR
        elif line.startswith('---'):
            html_parts.append('<hr>')

        # Blockquote
        elif line.startswith('> '):
            html_parts.append(f'<blockquote>{line[2:]}</blockquote>')

        # Bullet points
        elif line.startswith('- '):
            html_parts.append(f'<li>{line[2:]}</li>')

        # Empty line
        elif not line.strip():
            html_parts.append('<br>')

        # Regular text
        else:
            text = line.strip()
            if text:
                # Handle bold
                text = text.replace('**', '<b>').replace('**', '</b>')
                html_parts.append(f'<p>{text}</p>')

        i += 1

    if in_table:
        html_parts.append('</tbody></table>')

    return '\n'.join(html_parts)


def generate_pdf(html_content):
    """使用 WeasyPrint 生成 PDF"""
    with open(HTML_OUTPUT, 'w') as f:
        f.write(html_content)

    result = subprocess.run(
        ['weasyprint', HTML_OUTPUT, PDF_OUTPUT],
        capture_output=True, text=True, timeout=30
    )
    if result.returncode == 0:
        size = os.path.getsize(PDF_OUTPUT)
        print(f"PDF 生成成功: {PDF_OUTPUT} ({size//1024}KB)")
        return True
    else:
        print(f"PDF 生成失败: {result.stderr}")
        return False


def generate_ppt():
    """使用 python-pptx 生成 PPT"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme (BMW blue/white professional)
    DARK_BLUE = RGBColor(0x16, 0x21, 0x3e)
    ACCENT_RED = RGBColor(0xe9, 0x45, 0x60)
    MEDIUM_BLUE = RGBColor(0x0f, 0x34, 0x60)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

    def add_bg(slide, color=DARK_BLUE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=14,
                     bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        return tf

    def add_table_slide(slide, title, headers, rows, top=1.2):
        add_text_box(slide, 0.5, 0.3, 12, 0.6, title, 24, True, DARK_BLUE)
        n_rows = len(rows) + 1
        n_cols = len(headers)
        table_shape = slide.shapes.add_table(n_rows, n_cols,
            Inches(0.5), Inches(top), Inches(12.3), Inches(0.35 * n_rows))
        table = table_shape.table

        # Header
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = WHITE
                p.font.bold = True

        # Data rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i+1, j)
                cell.text = str(val)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = DARK_GRAY

    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BLUE)
    add_text_box(slide, 1, 2.5, 11, 1, "终端买家情报报告", 40, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, 2, 3.5, 9, 0.6, "澳大利亚 & 新西兰", 28, False, ACCENT_RED, PP_ALIGN.CENTER)
    add_text_box(slide, 2, 4.3, 9, 0.5, "钢结构折叠集成房屋 / 模块化建筑", 18, False, RGBColor(0xAA,0xBB,0xCC), PP_ALIGN.CENTER)
    add_text_box(slide, 3, 5.5, 7, 0.4, "2026-05-20 | 太一买家情报引擎", 12, False, RGBColor(0x88,0x99,0xAA), PP_ALIGN.CENTER)

    # Slide 2: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 0.5, 0.3, 12, 0.6, "报告摘要", 28, True, DARK_BLUE)
    summary_text = (
        "本报告覆盖澳大利亚 & 新西兰市场，共 10 家终端购买方/采购方\n\n"
        "■ 矿业营地模块化房屋 — 3家（BHP, Rio Tinto, Fortescue）\n"
        "■ 政府保障房/模块化住宅 — 2家（Kainga Ora, AusTender）\n"
        "■ 酒店集团模块化客房 — 5家（Accor, Hilton, Meriton, NZ Hotel, NZ Horizon）"
    )
    add_text_box(slide, 0.5, 1.3, 12, 3, summary_text, 16, False, DARK_GRAY)

    # Slides 3-12: Each buyer
    import json as _json
    buyer_data = [
        ("BHP Billiton（必和拓）", "矿业·澳大利亚", "矿工营房/营地模块",
         [["Luke King", "Head of Procurement", "luke.king@bhp.com"],
          ["Kurt Benavides", "Head of Procurement Operations", "kurt.benavides@bhp.com"],
          ["Tajinder Bedi", "Senior Procurement Manager", "LinkedIn"]]),

        ("Rio Tinto（力拓集团）", "矿业·澳大利亚", "矿工营房/营地设施",
         [["采购/供应链部", "供应商门户", "riotinto.com/en/suppliers"]]),

        ("Fortescue Metals Group", "矿业·澳大利亚", "矿工营房/营地设施",
         [["Dara Byrne", "Group Manager Contracts & Procurement", "dara.byrne@fmgl.com.au"],
          ["Mark Cocks", "Senior Contracts & Procurement Specialist", "LinkedIn"]]),

        ("Kainga Ora（政府住房署）", "政府·新西兰", "保障房/模块化住宅",
         [["Andrea Morton", "Director Procurement", "andrea.morton@kaingaora.govt.nz"],
          ["招标平台", "Tenderlink", "tenderlink.com/kaingaora"]]),

        ("AusTender（联邦政府采购）", "政府·澳大利亚", "预制建筑/模块化住宿",
         [["招标平台", "AusTender", "austender.gov.au"],
          ["关键词", "prefabricated building", "modular accommodation"]]),

        ("Accor Pacific（雅高太平洋）", "酒店·澳洲/NZ", "模块化酒店/度假屋",
         [["开发部/采购部", "LinkedIn搜索", "Accor Pacific Development Manager"]]),

        ("Hilton Australia（希尔顿澳洲）", "酒店·澳大利亚", "模块化酒店客房",
         [["已验证", "Townsville首个模块化酒店", "190间客房"]]),

        ("Meriton Group（美利通集团）", "酒店·澳大利亚", "酒店客房/服务公寓",
         [["规模", "6,211间客房", "澳洲最大酒店业主"]]),

        ("NZ Hotel Holdings", "酒店·新西兰", "酒店建设/模块化住宿",
         [["背景", "NZ Super Fund投资", "nzhotelholdings.co.nz"]]),

        ("NZ Horizon Hospitality Group", "酒店·新西兰南岛", "酒店客房模块",
         [["团队信息", "nzhhg.co.nz", ""]]),
    ]

    for name, industry, need, contacts in buyer_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text_box(slide, 0.5, 0.3, 12, 0.6, name, 24, True, DARK_BLUE)
        add_text_box(slide, 0.5, 0.9, 12, 0.4, f"{industry} | 采购需求: {need}", 14, False, ACCENT_RED)
        headers = ["姓名", "职位", "联系方式"]
        add_table_slide(slide, "关键联系人", headers, contacts, top=1.6)

    # Last slide: Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 0.5, 0.3, 12, 0.6, "触达策略建议", 28, True, DARK_BLUE)
    strategy = (
        "🔴 高优先级 — BHP, Fortescue, Kainga Ora\n"
        "    供应商门户注册 + LinkedIn 直联采购负责人\n\n"
        "🟡 中优先级 — Rio Tinto, AusTender, Hilton\n"
        "    招标平台订阅 + 开发信\n\n"
        "🟢 低优先级 — Accor, Meriton, NZ Hotel\n"
        "    长期关系建立\n\n"
        "推荐渠道：供应商门户注册 | LinkedIn 直连 | 招标平台订阅 | 行业协会"
    )
    add_text_box(slide, 0.5, 1.3, 12, 4, strategy, 14, False, DARK_GRAY)

    # Footer on all slides
    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "太一跨境贸易 Agent · 买家情报引擎", 9, False, RGBColor(0x99,0x99,0x99), PP_ALIGN.CENTER)

    prs.save(PPT_OUTPUT)
    size = os.path.getsize(PPT_OUTPUT)
    print(f"PPT 生成成功: {PPT_OUTPUT} ({size//1024}KB)")


def send_to_feishu_group(file_path, group_chat_id):
    """通过 lark-cli 发送文件到飞书群"""
    import subprocess
    env = os.environ.copy()
    env["LARK_CLI_NO_PROXY"] = "1"
    env["OPENCLAW_HOME"] = "/home/sayelf/.openclaw"

    # Get relative path from skill dir
    skill_dir = "/home/sayelf/openclaw/skills/轨道运营-agent"
    rel_path = os.path.relpath(file_path, skill_dir)

    result = subprocess.run(
        ["lark-cli", "--as", "bot", "im", "+messages-send",
         "--chat-id", group_chat_id, "--file", rel_path],
        capture_output=True, text=True, timeout=60,
        env=env, cwd=skill_dir
    )
    ok = result.returncode == 0 and '"ok": true' in result.stdout
    print(f"发送到飞书群: {os.path.basename(file_path)} {'✅' if ok else '❌'}")
    return ok


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    print("生成买家情报报告...")

    # Generate PPT
    print("\n--- 生成 PPT ---")
    try:
        generate_ppt()
    except ImportError:
        print("python-pptx 未安装，跳过 PPT")
    except Exception as e:
        print(f"PPT 生成失败: {e}")

    # Generate PDF
    print("\n--- 生成 PDF ---")
    html = md_to_html()
    try:
        generate_pdf(html)
    except Exception as e:
        print(f"PDF 生成失败: {e}")

    print("\n完成!")


if __name__ == "__main__":
    main()
